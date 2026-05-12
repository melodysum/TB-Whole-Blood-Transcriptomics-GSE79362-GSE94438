import sys
sys.path.insert(0, "/private/tmp/python_pkgs")

from pathlib import Path
import csv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

OUT = Path("/private/tmp/tb_curated_results")
FIG = OUT / "figures"
PPTX = OUT / "TB_curated_transcriptomics_GSE79362_GSE94438.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

COLORS = {
    "ink": RGBColor(31, 41, 55),
    "muted": RGBColor(107, 114, 128),
    "bg": RGBColor(247, 250, 252),
    "teal": RGBColor(0, 121, 128),
    "red": RGBColor(178, 34, 52),
    "orange": RGBColor(214, 118, 34),
    "green": RGBColor(45, 126, 82),
    "blue": RGBColor(46, 93, 160),
    "light": RGBColor(230, 236, 240),
}

def read_csv(name):
    with open(OUT / name, newline="") as f:
        return list(csv.DictReader(f))

summary = {r["study"]: r for r in read_csv("analysis_summary.csv")}
sig = read_csv("signature_AUC_summary.csv")
cross = {r["metric"]: r["value"] for r in read_csv("cross_dataset_summary.csv")}

def set_bg(slide, color=COLORS["bg"]):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = COLORS["ink"]
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.57), Inches(0.88), Inches(12.1), Inches(0.35))
        q = sub.text_frame.paragraphs[0]
        q.text = subtitle
        q.font.size = Pt(11)
        q.font.color.rgb = COLORS["muted"]

def add_footer(slide, text="curatedTBData hg38 reprocessed counts; baseline edgeR unless noted"):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.12), Inches(12.2), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(8.5)
    p.font.color.rgb = COLORS["muted"]

def add_bullets(slide, items, x, y, w, h, size=14, color=COLORS["ink"]):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0
        p.space_after = Pt(6)
    return box

def add_card(slide, x, y, w, h, headline, body, color=COLORS["teal"]):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = COLORS["light"]
    shape.adjustments[0] = 0.06
    t = slide.shapes.add_textbox(Inches(x+0.18), Inches(y+0.14), Inches(w-0.36), Inches(0.35))
    p = t.text_frame.paragraphs[0]
    p.text = headline
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    b = slide.shapes.add_textbox(Inches(x+0.18), Inches(y+0.55), Inches(w-0.36), Inches(h-0.65))
    q = b.text_frame.paragraphs[0]
    q.text = body
    q.font.size = Pt(11)
    q.font.color.rgb = COLORS["ink"]

def add_image(slide, path, x, y, w, h=None):
    if Path(path).exists():
        if h:
            slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        else:
            slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))

def add_table(slide, rows, x, y, w, h, font_size=10):
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = COLORS["ink"]
                if r == 0:
                    p.font.bold = True
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["light"]
    return table

def slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    return s

# 1
s = slide()
set_bg(s, RGBColor(235, 244, 246))
add_title(s, "TB whole-blood transcriptomics: infection, progression, and disease signals",
          "GSE79362 PTB/LTBI and GSE94438 household-contact progression cohort")
add_card(s, 0.7, 1.45, 3.7, 1.35, "核心问题", "哪些 whole-blood RNA 信号反映结核发病，而不是单纯暴露或潜伏感染？", COLORS["teal"])
add_card(s, 4.8, 1.45, 3.7, 1.35, "关键比较", "GSE79362: LTBI vs PTB；GSE94438: household-contact Control vs PTB/progressor。", COLORS["blue"])
add_card(s, 8.9, 1.45, 3.7, 1.35, "主要结论", "两个队列共同指向 IFN、炎症、补体、髓系/中性粒细胞相关通路。", COLORS["red"])
add_bullets(s, [
    "这里的“未感染”需要谨慎：GSE94438 Control 是密切接触者/未进展者，不等价于严格未感染。",
    "GSE79362 更适合回答：在已感染背景下，哪些信号与活动性 TB 相关。",
    "GSE94438 更适合回答：在暴露密接人群中，哪些信号与后续 TB 进展/发病相关。"
], 0.95, 3.45, 11.4, 2.0, 15)
add_footer(s, "Prepared from actual curatedTBData analysis outputs in /private/tmp/tb_curated_results")

# 2
s = slide()
add_title(s, "Dataset structure: the two cohorts ask related but not identical biological questions")
rows = [
    ["Dataset", "Biological contrast", "Labelled n", "Site/country", "Repeated sampling"],
    ["GSE79362", "PTB vs LTBI", "355 = 110 PTB + 245 LTBI", "South Africa only", "Yes: 144 patients; 105 repeated"],
    ["GSE94438", "PTB/progressor vs household-contact Control", "428 labelled = 101 PTB + 327 Control", "Ethiopia / South Africa / The Gambia", "Yes: 334 patients; 79 repeated"],
]
add_table(s, rows, 0.6, 1.3, 12.1, 1.45, 10)
add_card(s, 0.8, 3.15, 3.7, 1.35, "GSE79362 解释重点", "全体为 South Africa；PTB/LTBI 让我们看见“潜伏感染背景上的活动性疾病信号”。", COLORS["teal"])
add_card(s, 4.85, 3.15, 3.7, 1.35, "GSE94438 解释重点", "三国密接队列；site 是优先混杂因素，Control 不是严格健康未暴露人群。", COLORS["orange"])
add_card(s, 8.9, 3.15, 3.7, 1.35, "分析限制", "重复采样和采样时间可能影响有效样本量；差异结果需看调整模型敏感性。", COLORS["red"])
add_bullets(s, [
    "GSE94438 GEO total 为 434；curated hg38 也有 434，但 6 个样本缺 TBStatus，因此 labelled analysis 为 428。",
    "GSE79362 curated hg38 保留 355/355。"
], 0.85, 5.1, 11.6, 0.9, 13)
add_footer(s)

# 3
s = slide()
add_title(s, "QC: sample-level variation is visible and must be separated from biology")
add_image(s, FIG / "GSE79362_PCA_TBStatus.png", 0.55, 1.15, 5.85, 4.8)
add_image(s, FIG / "GSE94438_PCA_site.png", 6.8, 1.15, 5.85, 4.8)
add_bullets(s, [
    "PCA by TBStatus helps assess whether disease/progression is a major source of variance.",
    "GSE94438 must be checked by country/site because geography and population structure can mimic biology.",
    "Library-size and correlation heatmaps were generated as QC outputs; no low-library-size filtering threshold was applied beyond gene-level CPM filtering."
], 0.75, 6.1, 11.8, 0.8, 10.8)
add_footer(s, "QC plots: library size, PCA by TBStatus, PCA by site/country, sample correlation heatmaps")

# 4
s = slide()
add_title(s, "Differential expression: activity/progression is dominated by up-regulated innate immune genes")
rows = [
    ["Dataset", "Baseline DEGs", "Up", "Down", "Adjusted sensitivity"],
    ["GSE79362", "30", "29", "1", "9 after timepoint + PatientID block"],
    ["GSE94438", "43", "43", "0", "41 after site + sex + age"],
]
add_table(s, rows, 0.65, 1.1, 5.0, 1.3, 11)
add_image(s, FIG / "GSE79362_volcano.png", 0.55, 2.75, 5.7, 3.9)
add_image(s, FIG / "GSE94438_volcano.png", 6.7, 2.75, 5.7, 3.9)
add_card(s, 6.35, 1.1, 6.1, 1.3, "Top disease-associated genes", "GSE79362: GBP6, FCGR1CP, GBP5, ANKRD22, PDCD1LG2, CD274, FCGR1B.  GSE94438: SEPTIN4, ANKRD22, C1QC, BATF2, SERPING1, C1QB, CD274.", COLORS["red"])
add_footer(s)

# 5
s = slide()
add_title(s, "Adjustment matters differently in the two cohorts")
add_image(s, FIG / "GSE79362_baseline_vs_adjusted.png", 0.65, 1.15, 5.75, 4.7)
add_image(s, FIG / "GSE94438_baseline_vs_adjusted.png", 6.85, 1.15, 5.75, 4.7)
add_bullets(s, [
    "GSE79362: repeated longitudinal structure is non-trivial; duplicateCorrelation estimated within-patient correlation ≈ 0.306, and strict DEGs dropped from 30 to 9.",
    "GSE94438: site/sex/age adjustment preserved most disease signal, with strict DEGs 43 to 41.",
    "Interpretation: GSE79362 disease signal is real but sensitive to repeated sampling/time; GSE94438 signal is not simply explained by site."
], 0.8, 6.05, 11.8, 1.0, 11)
add_footer(s)

# 6
s = slide()
add_title(s, "Pathway enrichment: both datasets converge on IFN, inflammation, complement, and innate immunity")
add_image(s, FIG / "GSE79362_GSEA_Hallmark.png", 0.55, 1.05, 5.95, 4.75)
add_image(s, FIG / "GSE94438_GSEA_Hallmark.png", 6.75, 1.05, 5.95, 4.75)
add_bullets(s, [
    "Hallmark IFN-α, IFN-γ, inflammatory response, TNF/NFκB, IL6/JAK/STAT3, and complement are strongly positive in both datasets.",
    "These pathways are canonical TB whole-blood signals, but they may reflect leukocyte composition shifts as well as within-cell activation."
], 0.8, 6.05, 11.8, 0.75, 11)
add_footer(s, "GSEA ranking: sign(logFC) × -log10(PValue); FDR values for key Hallmark pathways are extremely small")

# 7
s = slide()
add_title(s, "Signature validation: risk scores work better in PTB/LTBI than in household-contact progression")
rows = [["Signature", "GSE79362 AUC", "GSE94438 AUC", "Missing genes"]]
for name in ["Zak16", "RISK4", "Eleven_gene"]:
    r1 = next(r for r in sig if r["study"] == "GSE79362" and r["Signature"] == name)
    r2 = next(r for r in sig if r["study"] == "GSE94438" and r["Signature"] == name)
    miss = r2["missing_genes"] if r2["missing_genes"] else "none"
    rows.append([name, f'{float(r1["AUC"]):.3f}', f'{float(r2["AUC"]):.3f}', miss])
add_table(s, rows, 0.7, 1.05, 5.35, 1.65, 11)
add_image(s, FIG / "GSE79362_Zak16_score.png", 0.65, 3.0, 5.35, 3.4)
add_image(s, FIG / "GSE94438_Zak16_score_by_site.png", 6.35, 2.0, 6.35, 4.4)
add_bullets(s, [
    "AUC is consistently higher in GSE79362 (≈0.76-0.77) than GSE94438 (≈0.69).",
    "Likely explanation: PTB/LTBI contrast is biologically sharper than predicting progression among heterogeneous household contacts.",
    "GSE94438 site-stratified plots are essential because score distributions vary by country."
], 6.35, 6.35, 6.1, 0.8, 10)
add_footer(s)

# 8
s = slide()
add_title(s, "Cross-dataset consistency: pathways agree more strongly than individual genes")
add_image(s, FIG / "cross_dataset_logFC_scatter.png", 0.65, 1.15, 5.7, 4.7)
add_image(s, FIG / "cross_dataset_Hallmark_NES_scatter.png", 6.9, 1.15, 5.55, 4.7)
add_card(s, 0.85, 6.05, 3.3, 0.75, "Gene-level", f"Shared tested genes: {int(float(cross['shared_genes_tested'])):,}; logFC Spearman = {float(cross['spearman_logFC']):.2f}", COLORS["blue"])
add_card(s, 4.55, 6.05, 3.3, 0.75, "Strict DEG overlap", f"Shared strict DEGs: {int(float(cross['shared_sig_DEGs']))}; top-100 overlap: {int(float(cross['top100_overlap']))}", COLORS["orange"])
add_card(s, 8.25, 6.05, 3.3, 0.75, "Pathway-level", f"Hallmark NES Spearman = {float(cross['hallmark_NES_spearman']):.2f}", COLORS["green"])
add_footer(s)

# 9
s = slide()
add_title(s, "Biological interpretation: what is common to infection, disease, and non-progression?")
add_card(s, 0.75, 1.25, 3.7, 1.55, "Infected but not diseased / non-progressor", "LTBI or household-contact Control samples show lower IFN/myeloid/complement score on average; they are not transcriptionally inert, but lack the strong inflammatory disease program.", COLORS["green"])
add_card(s, 4.8, 1.25, 3.7, 1.55, "Progression / active TB", "Common signal: IFN-inducible GBP/ISG genes, antigen/inhibitory checkpoint genes such as CD274, complement C1 genes, and myeloid activation.", COLORS["red"])
add_card(s, 8.85, 1.25, 3.7, 1.55, "Strictly uninfected", "Not directly resolved here. GSE94438 controls are exposed household contacts; TST/QFT metadata would be needed to define uninfected subgroups.", COLORS["orange"])
add_bullets(s, [
    "共性一：发病/进展状态在两个队列都呈现强 IFN 和 innate immune activation。",
    "共性二：补体和髓系/中性粒细胞通路说明 whole blood 信号很可能受细胞组成影响。",
    "共性三：签名基因表现为连续风险梯度，而不是完美二分类；这符合 TB 从暴露、感染、亚临床到活动性疾病的连续谱。"
], 0.9, 3.55, 11.8, 2.0, 15)
add_footer(s, "Do not overinterpret whole-blood DEGs as cell-intrinsic regulation without deconvolution or single-cell validation")

# 10
s = slide()
add_title(s, "Limitations and next analyses")
add_bullets(s, [
    "curatedTBData subset: GSE79362 retains 355/355; GSE94438 retains 434 hg38 counts but 6 samples lack TBStatus, leaving 428 labelled samples.",
    "Label definitions differ: GSE79362 is PTB/LTBI, while GSE94438 is household-contact Control/PTB; this affects AUC and DEG interpretation.",
    "Repeated measures matter: GSE79362 adjustment reduced DEGs substantially; baseline-only sensitivity analysis should be considered.",
    "Site/population effect matters most for GSE94438; site adjustment preserved most DEGs but should remain in the final model.",
    "Whole-blood composition is a major caveat; next step should include cell-type deconvolution or validation in sorted/single-cell data."
], 0.85, 1.2, 11.9, 3.3, 15)
add_card(s, 1.0, 5.15, 11.25, 1.0, "Take-home message", "Across two prospective TB whole-blood RNA-seq cohorts, active/progressive TB is marked by a reproducible IFN–inflammatory–complement–myeloid program; this signal distinguishes disease risk from LTBI/exposed non-progression, but it likely mixes immune activation with leukocyte composition changes.", COLORS["teal"])
add_footer(s)

prs.save(PPTX)
print(PPTX)
