import sys
sys.path.insert(0, "/private/tmp/python_pkgs")

from pathlib import Path
import csv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

OUT = Path("/private/tmp/tb_curated_results")
FIG = OUT / "figures"
PPTX = OUT / "TB_curated_transcriptomics_GSE79362_GSE94438_English.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

COLORS = {
    "ink": RGBColor(31, 41, 55),
    "muted": RGBColor(107, 114, 128),
    "bg": RGBColor(247, 250, 252),
    "teal": RGBColor(0, 121, 128),
    "red": RGBColor(178, 34, 52),
    "orange": RGBColor(214, 118, 34),
    "green": RGBColor(45, 126, 82),
    "blue": RGBColor(46, 93, 160),
    "light": RGBColor(230, 236, 240),
}

def read_csv(name):
    with open(OUT / name, newline="") as f:
        return list(csv.DictReader(f))

sig = read_csv("signature_AUC_summary.csv")
cross = {r["metric"]: r["value"] for r in read_csv("cross_dataset_summary.csv")}

def set_bg(slide, color=COLORS["bg"]):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = COLORS["ink"]
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.57), Inches(0.88), Inches(12.1), Inches(0.35))
        q = sub.text_frame.paragraphs[0]
        q.text = subtitle
        q.font.size = Pt(11)
        q.font.color.rgb = COLORS["muted"]

def add_footer(slide, text="curatedTBData hg38 reprocessed counts; baseline edgeR unless stated otherwise"):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.12), Inches(12.2), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(8.5)
    p.font.color.rgb = COLORS["muted"]

def add_bullets(slide, items, x, y, w, h, size=14, color=COLORS["ink"]):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0
        p.space_after = Pt(6)
    return box

def add_card(slide, x, y, w, h, headline, body, color=COLORS["teal"]):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = COLORS["light"]
    shape.adjustments[0] = 0.06
    t = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.14), Inches(w - 0.36), Inches(0.35))
    p = t.text_frame.paragraphs[0]
    p.text = headline
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    b = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.55), Inches(w - 0.36), Inches(h - 0.65))
    q = b.text_frame.paragraphs[0]
    q.text = body
    q.font.size = Pt(11)
    q.font.color.rgb = COLORS["ink"]

def add_image(slide, path, x, y, w, h=None):
    if Path(path).exists():
        if h:
            slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        else:
            slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))

def add_table(slide, rows, x, y, w, h, font_size=10):
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = COLORS["ink"]
                if r == 0:
                    p.font.bold = True
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["light"]
    return table

def slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    return s

s = slide()
set_bg(s, RGBColor(235, 244, 246))
add_title(s, "TB whole-blood transcriptomics: infection, progression, and active disease",
          "GSE79362 PTB/LTBI and GSE94438 household-contact progression cohort")
add_card(s, 0.7, 1.45, 3.7, 1.35, "Central question", "Which whole-blood RNA signals reflect TB disease rather than exposure or latent infection alone?", COLORS["teal"])
add_card(s, 4.8, 1.45, 3.7, 1.35, "Key comparison", "GSE79362: LTBI versus PTB. GSE94438: household-contact Control versus PTB/progressor.", COLORS["blue"])
add_card(s, 8.9, 1.45, 3.7, 1.35, "Main finding", "Both cohorts converge on interferon, inflammation, complement, and myeloid/neutrophil-associated pathways.", COLORS["red"])
add_bullets(s, [
    "The term “uninfected” should be used cautiously: GSE94438 controls are exposed household contacts or non-progressors, not strictly unexposed healthy individuals.",
    "GSE79362 is best suited to identify active-disease signals on a latent-infection background.",
    "GSE94438 is best suited to identify progression/disease-emergence signals among exposed household contacts."
], 0.95, 3.45, 11.4, 2.0, 15)
add_footer(s, "Prepared from actual curatedTBData analysis outputs in /private/tmp/tb_curated_results")

s = slide()
add_title(s, "Dataset structure: the two cohorts address related but distinct biological contrasts")
rows = [
    ["Dataset", "Biological contrast", "Labelled n", "Site/country", "Repeated sampling"],
    ["GSE79362", "PTB vs LTBI", "355 = 110 PTB + 245 LTBI", "South Africa only", "Yes: 144 patients; 105 repeated"],
    ["GSE94438", "PTB/progressor vs household-contact Control", "428 labelled = 101 PTB + 327 Control", "Ethiopia / South Africa / The Gambia", "Yes: 334 patients; 79 repeated"],
]
add_table(s, rows, 0.6, 1.3, 12.1, 1.45, 10)
add_card(s, 0.8, 3.15, 3.7, 1.35, "How to interpret GSE79362", "All samples are from South Africa; PTB/LTBI primarily captures active-disease biology on a latent-infection background.", COLORS["teal"])
add_card(s, 4.85, 3.15, 3.7, 1.35, "How to interpret GSE94438", "This is a three-country household-contact cohort; site is a priority confounder and Control does not mean unexposed.", COLORS["orange"])
add_card(s, 8.9, 3.15, 3.7, 1.35, "Analytical caveat", "Repeated sampling and sampling time can inflate the effective sample size; adjusted sensitivity models are therefore essential.", COLORS["red"])
add_bullets(s, [
    "GSE94438 has 434 hg38 count samples in curatedTBData, but 6 lack TBStatus; labelled analyses therefore use 428 samples.",
    "GSE79362 retains 355/355 samples in the curated hg38 matrix."
], 0.85, 5.1, 11.6, 0.9, 13)
add_footer(s)

s = slide()
add_title(s, "QC: sample-level variation must be separated from biological signal")
add_image(s, FIG / "GSE79362_PCA_TBStatus.png", 0.55, 1.15, 5.85, 4.8)
add_image(s, FIG / "GSE94438_PCA_site.png", 6.8, 1.15, 5.85, 4.8)
add_bullets(s, [
    "PCA by TBStatus assesses whether active disease or progression is a major source of variance.",
    "GSE94438 must also be inspected by country/site because geography and population structure can mimic disease biology.",
    "Library-size histograms and sample-correlation heatmaps were generated; gene-level CPM filtering was applied before downstream analyses."
], 0.75, 6.1, 11.8, 0.8, 10.8)
add_footer(s, "QC outputs: library size, PCA by TBStatus, PCA by site/country, sample-correlation heatmaps")

s = slide()
add_title(s, "Differential expression: active/progressive TB is dominated by up-regulated innate immune genes")
rows = [
    ["Dataset", "Baseline DEGs", "Up", "Down", "Adjusted sensitivity"],
    ["GSE79362", "30", "29", "1", "9 after timepoint + PatientID block"],
    ["GSE94438", "43", "43", "0", "41 after site + sex + age"],
]
add_table(s, rows, 0.65, 1.1, 5.0, 1.3, 11)
add_image(s, FIG / "GSE79362_volcano.png", 0.55, 2.75, 5.7, 3.9)
add_image(s, FIG / "GSE94438_volcano.png", 6.7, 2.75, 5.7, 3.9)
add_card(s, 6.35, 1.1, 6.1, 1.3, "Top disease-associated genes", "GSE79362: GBP6, FCGR1CP, GBP5, ANKRD22, PDCD1LG2, CD274, FCGR1B. GSE94438: SEPTIN4, ANKRD22, C1QC, BATF2, SERPING1, C1QB, CD274.", COLORS["red"])
add_footer(s)

s = slide()
add_title(s, "Adjustment affects the two cohorts differently")
add_image(s, FIG / "GSE79362_baseline_vs_adjusted.png", 0.65, 1.15, 5.75, 4.7)
add_image(s, FIG / "GSE94438_baseline_vs_adjusted.png", 6.85, 1.15, 5.75, 4.7)
add_bullets(s, [
    "GSE79362 has a meaningful longitudinal structure: duplicateCorrelation estimated a within-patient correlation of approximately 0.306, and strict DEGs decreased from 30 to 9.",
    "GSE94438 retained most disease-associated signal after adjustment for site, sex, and age, with strict DEGs decreasing only from 43 to 41.",
    "Interpretation: the GSE79362 signal is sensitive to repeated sampling and timepoint, whereas the GSE94438 signal is not simply explained by site."
], 0.8, 6.05, 11.8, 1.0, 11)
add_footer(s)

s = slide()
add_title(s, "Pathway enrichment: both datasets converge on interferon, inflammation, complement, and innate immunity")
add_image(s, FIG / "GSE79362_GSEA_Hallmark.png", 0.55, 1.05, 5.95, 4.75)
add_image(s, FIG / "GSE94438_GSEA_Hallmark.png", 6.75, 1.05, 5.95, 4.75)
add_bullets(s, [
    "Hallmark interferon-alpha, interferon-gamma, inflammatory response, TNF/NF-kB, IL6/JAK/STAT3, and complement pathways are strongly positive in both datasets.",
    "These are canonical TB whole-blood signatures, but they may reflect leukocyte-composition shifts as well as cell-intrinsic activation."
], 0.8, 6.05, 11.8, 0.75, 11)
add_footer(s, "GSEA ranking: sign(logFC) x -log10(PValue); key Hallmark pathways show extremely small adjusted P values")

s = slide()
add_title(s, "Signature validation: risk scores perform better in PTB/LTBI than in household-contact progression")
rows = [["Signature", "GSE79362 AUC", "GSE94438 AUC", "Missing genes"]]
for name in ["Zak16", "RISK4", "Eleven_gene"]:
    r1 = next(r for r in sig if r["study"] == "GSE79362" and r["Signature"] == name)
    r2 = next(r for r in sig if r["study"] == "GSE94438" and r["Signature"] == name)
    miss = r2["missing_genes"] if r2["missing_genes"] else "none"
    rows.append([name, f'{float(r1["AUC"]):.3f}', f'{float(r2["AUC"]):.3f}', miss])
add_table(s, rows, 0.7, 1.05, 5.35, 1.65, 11)
add_image(s, FIG / "GSE79362_Zak16_score.png", 0.65, 3.0, 5.35, 3.4)
add_image(s, FIG / "GSE94438_Zak16_score_by_site.png", 6.35, 2.0, 6.35, 4.4)
add_bullets(s, [
    "AUCs are consistently higher in GSE79362, approximately 0.76-0.77, than in GSE94438, approximately 0.69.",
    "This likely reflects a sharper PTB/LTBI contrast compared with the harder task of distinguishing progressors from heterogeneous household contacts.",
    "For GSE94438, site-stratified score plots are essential because score distributions vary by country."
], 6.35, 6.35, 6.1, 0.8, 10)
add_footer(s)

s = slide()
add_title(s, "Cross-dataset consistency: pathways agree more strongly than individual genes")
add_image(s, FIG / "cross_dataset_logFC_scatter.png", 0.65, 1.15, 5.7, 4.7)
add_image(s, FIG / "cross_dataset_Hallmark_NES_scatter.png", 6.9, 1.15, 5.55, 4.7)
add_card(s, 0.85, 6.05, 3.3, 0.75, "Gene level", f"Shared tested genes: {int(float(cross['shared_genes_tested'])):,}; logFC Spearman = {float(cross['spearman_logFC']):.2f}", COLORS["blue"])
add_card(s, 4.55, 6.05, 3.3, 0.75, "Strict DEG overlap", f"Shared strict DEGs: {int(float(cross['shared_sig_DEGs']))}; top-100 overlap: {int(float(cross['top100_overlap']))}", COLORS["orange"])
add_card(s, 8.25, 6.05, 3.3, 0.75, "Pathway level", f"Hallmark NES Spearman = {float(cross['hallmark_NES_spearman']):.2f}", COLORS["green"])
add_footer(s)

s = slide()
add_title(s, "Biological interpretation: what is shared by infection, disease, and non-progression?")
add_card(s, 0.75, 1.25, 3.7, 1.55, "Infected but not diseased / non-progressor", "LTBI and household-contact Control samples show lower interferon, myeloid, and complement scores on average; they are not transcriptionally inert, but they lack the strong inflammatory disease program.", COLORS["green"])
add_card(s, 4.8, 1.25, 3.7, 1.55, "Progression / active TB", "The shared signal includes interferon-inducible GBP/ISG genes, immune-regulatory genes such as CD274, complement C1 genes, and myeloid activation.", COLORS["red"])
add_card(s, 8.85, 1.25, 3.7, 1.55, "Strictly uninfected individuals", "This category is not directly resolved here. GSE94438 controls are exposed household contacts; TST/QFT metadata would be required to define uninfected subgroups.", COLORS["orange"])
add_bullets(s, [
    "Commonality 1: active or progressive TB shows strong interferon and innate immune activation in both cohorts.",
    "Commonality 2: complement and myeloid/neutrophil pathway enrichment suggests that whole-blood signals are likely influenced by leukocyte composition.",
    "Commonality 3: signature scores behave as a continuous risk gradient rather than a perfect binary classifier, consistent with the biological continuum from exposure to infection, subclinical disease, and active TB."
], 0.9, 3.55, 11.8, 2.0, 15)
add_footer(s, "Whole-blood DEGs should not be interpreted as cell-intrinsic regulation without deconvolution or single-cell validation")

s = slide()
add_title(s, "Limitations and next analyses")
add_bullets(s, [
    "curatedTBData subset: GSE79362 retains 355/355 samples; GSE94438 retains 434 hg38 count samples but 6 lack TBStatus, leaving 428 labelled samples.",
    "Label definitions differ: GSE79362 is PTB/LTBI, whereas GSE94438 is household-contact Control/PTB; this affects both AUC and DEG interpretation.",
    "Repeated measures matter: adjustment for timepoint and patient blocking substantially reduced GSE79362 DEGs; a baseline-only sensitivity analysis should also be considered.",
    "Site and population effects are especially important for GSE94438; site adjustment preserved most DEGs but should remain in the final model.",
    "Whole-blood composition is a major caveat; the next step should include cell-type deconvolution or validation in sorted-cell/single-cell data."
], 0.85, 1.2, 11.9, 3.3, 15)
add_card(s, 1.0, 5.15, 11.25, 1.0, "Take-home message", "Across two prospective TB whole-blood RNA-seq cohorts, active or progressive TB is marked by a reproducible interferon-inflammatory-complement-myeloid program; this signal distinguishes disease risk from LTBI or exposed non-progression, but it likely combines immune activation with leukocyte-composition changes.", COLORS["teal"])
add_footer(s)

prs.save(PPTX)
print(PPTX)
